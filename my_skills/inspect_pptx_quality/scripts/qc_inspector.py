import sys
import os
import json
from pptx import Presentation

def inspect_pptx(pptx_path, brief_path):
    print(f"--- QC Inspection ---")
    print(f"Target: {pptx_path}")
    
    # 1. Existence
    if not os.path.exists(pptx_path):
        print("FAIL: PPTX file not found.")
        sys.exit(1)
        
    if not os.path.exists(brief_path):
        print("FAIL: Brief file not found.")
        sys.exit(1)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"FAIL: Invalid PPTX file. {e}")
        sys.exit(1)
        
    with open(brief_path, 'r') as f:
        brief = json.load(f)
        
    required_layouts = brief.get('required_layouts', [])
    expected_count = len(required_layouts)
    actual_count = len(prs.slides)
    
    # 2. Slide Count
    if expected_count != actual_count:
        print(f"FAIL: Slide count mismatch. Expected {expected_count}, got {actual_count}.")
        sys.exit(1)
    else:
        print(f"PASS: Slide count matches ({actual_count}).")
        
    # 3. Content Sampling
    print("Sampling content...")
    failures = []
    
    for i, slide in enumerate(prs.slides):
        req = required_layouts[i]
        
        # Check Background (At least 1 pic)
        pics = [shape for shape in slide.shapes if shape.shape_type == 13] # MSO_SHAPE_TYPE.PICTURE = 13
        if len(pics) < 1:
            failures.append(f"Slide {i+1}: Missing background image.")
            
        # Check Title (if brief has title)
        if 'title' in req:
            expected_title = req['title']
            found_title = False
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                # Simple loose matching
                if expected_title in shape.text_frame.text:
                    found_title = True
                    break
            
            if not found_title:
                failures.append(f"Slide {i+1}: Expected title '{expected_title}' not found.")

    if len(failures) > 0:
        print("FAIL: Quality issues found:")
        for fail in failures:
            print(f"  - {fail}")
        sys.exit(1)

    print("PASS: All checks passed.")
    sys.exit(0)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python qc_inspector.py <pptx_path> <brief_path>")
        sys.exit(1)
        
    inspect_pptx(sys.argv[1], sys.argv[2])
