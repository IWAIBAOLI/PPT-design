
import json
import re
from pathlib import Path
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

def parse_rgb(rgb_string):
    """Parse 'rgb(255, 0, 255)' to RGBColor."""
    match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', rgb_string)
    if match:
        return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))
    return RGBColor(255, 255, 255)  # Default white

def px_to_inches(px, dpi=96):
    """Convert pixels to inches."""
    return px / dpi

def create_pptx(screenshots_dir, output_file):
    """Create PPTX from screenshots with pure visual approach."""
    # Initialize presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 = 10:5.625
    prs.slide_height = Inches(5.625)
    
    screenshots_path = Path(screenshots_dir)
    
    # Get all PNG screenshot files in order
    screenshot_files = sorted(list(screenshots_path.glob("slide_*.png")))
    print(f"== Building PPTX with {len(screenshot_files)} slides ==")
    
    if len(screenshot_files) == 0:
        print(f"!! No PNG files found in {screenshots_path}")
        print(f"   Please run capture_screenshots.py first")
        return
    
    for screenshot_file in screenshot_files:
        print(f"-- Processing {screenshot_file.stem} --")
        
        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Insert background screenshot as full-slide image
        slide.shapes.add_picture(
            str(screenshot_file),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"   -> Added: {screenshot_file.name}")
    
    # Save presentation
    prs.save(output_file)
    print(f"\n✅ PPTX saved to: {output_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Create PPTX from screenshots")
    parser.add_argument("screenshots_dir", help="Directory containing PNG screenshots")
    parser.add_argument("output_file", help="Output .pptx file path")
    args = parser.parse_args()
    
    create_pptx(args.screenshots_dir, args.output_file)
