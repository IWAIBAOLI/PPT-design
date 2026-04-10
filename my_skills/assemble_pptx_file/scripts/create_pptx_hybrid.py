
import json
from pathlib import Path
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_color(color_str):
    """Parse CSS color to RGBColor."""
    if color_str.startswith('rgb'):
        match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if match:
            return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))
    elif color_str.startswith('#'):
        # Handle hex colors
        hex_color = color_str.lstrip('#')
        if len(hex_color) == 6:
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    return RGBColor(26, 26, 26)  # Default dark color

def px_to_emu(px):
    """Convert pixels to EMUs (English Metric Units)."""
    # 1 inch = 914400 EMUs, 1 inch = 96 px (standard DPI)
    return int(px * 914400 / 96)

def create_pptx_hybrid(screenshots_dir, metrics_dir, output_file):
    """
    Create PPTX using Hybrid Assembly Strategy:
    - Background: Screenshot with visual effects (text hidden)
    - Foreground: Editable TextBoxes for all text content
    """
    # Initialize presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 = 10:5.625
    prs.slide_height = Inches(5.625)
    
    screenshots_path = Path(screenshots_dir)
    metrics_path = Path(metrics_dir)
    
    # Get all screenshot files
    screenshot_files = sorted(list(screenshots_path.glob("slide_*.png")))
    print(f"== Building HYBRID PPTX with {len(screenshot_files)} slides ==")
    
    if len(screenshot_files) == 0:
        print(f"!! No PNG files found in {screenshots_path}")
        return
    
    for screenshot_file in screenshot_files:
        slide_name = screenshot_file.stem
        print(f"-- Processing {slide_name} --")
        
        # Load metrics
        metrics_file = metrics_path / f"{slide_name}.json"
        if not metrics_file.exists():
            print(f"   !! Warning: No metrics found for {slide_name}")
            continue
            
        with open(metrics_file, 'r', encoding='utf-8') as f:
            metrics_data = json.load(f)
        
        # Metrics JSON structure: {"elements": [{...}, {...}]}
        elements = metrics_data.get('elements', [])
        
        # Add blank slide
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Layer 1: Insert background screenshot (with hidden text)
        slide.shapes.add_picture(
            str(screenshot_file),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"   -> Background: {screenshot_file.name}")
        
        # Layer 2: Add editable text boxes based on metrics
        text_elements_added = 0
        
        for elem in elements:
            elem_role = elem.get('role', '')
            text_content = elem.get('text', '').strip()  # Use 'text' not 'text_content'
            
            if not text_content:
                continue
            
            bbox = elem.get('rect', {})
            styles = elem.get('style', {})
            
            # Get font size for width calculation
            font_size_str = styles.get('fontSize', '16px')
            font_size_px = float(font_size_str.replace('px', ''))
            
            # CRITICAL FIX: Calculate proper width based on text length
            # Metrics width is often wrong (single character width)
            # Use heuristic: character_count * font_size * 0.6 (average char width ratio)
            char_count = len(text_content)
            estimated_width_px = max(char_count * font_size_px * 0.6, 100)  # Minimum 100px
            
            # For Title elements, use larger width (often full slide width)
            if elem_role == 'Title' or 'title' in elem.get('selector', '').lower():
                estimated_width_px = 1600  # Use most of slide width for titles
            
            # Create text box with calculated dimensions
            left_emu = px_to_emu(bbox.get('x', 0))
            top_emu = px_to_emu(bbox.get('y', 0))
            width_emu = px_to_emu(estimated_width_px)  # Use calculated width
            height_emu = px_to_emu(max(font_size_px * 1.5, 50))  # Height based on font size
            
            try:
                txBox = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
                tf = txBox.text_frame
                tf.text = text_content
                tf.word_wrap = True
                
                # Apply formatting
                if tf.paragraphs:
                    para = tf.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        
                        # Font size (already calculated above)
                        font_size_pt = max(8, min(72, font_size_px * 0.75))
                        run.font.size = Pt(font_size_pt)
                        
                        # Font color
                        color_str = styles.get('color', '#1A1A1A')
                        run.font.color.rgb = parse_color(color_str)
                        
                        # Font family
                        font_family = styles.get('fontFamily', 'Arial')
                        run.font.name = font_family.split(',')[0].strip().replace('"', '')
                        
                        # Font weight
                        font_weight = styles.get('fontWeight', '400')
                        if int(font_weight) >= 600:
                            run.font.bold = True
                
                text_elements_added += 1
                
            except Exception as e:
                print(f"   !! Error creating text box for {elem_role}: {e}")
        
        print(f"   -> Added {text_elements_added} editable text boxes")
    
    # Save presentation
    prs.save(output_file)
    print(f"\n✅ HYBRID PPTX saved to: {output_file}")
    print(f"   Text boxes are EDITABLE, visual effects preserved in background")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Create Hybrid PPTX (editable text + visual background)")
    parser.add_argument("screenshots_dir", help="Directory containing PNG screenshots (with hidden text)")
    parser.add_argument("metrics_dir", help="Directory containing JSON metrics")
    parser.add_argument("output_file", help="Output .pptx file path")
    args = parser.parse_args()
    
    create_pptx_hybrid(args.screenshots_dir, args.metrics_dir, args.output_file)
