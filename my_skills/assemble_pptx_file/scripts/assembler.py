import os
import json
import sys
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def css_rgb_to_tuple(css_rgb):
    # css_rgb format: "rgb(255, 255, 255)"
    parts = css_rgb.replace('rgb(', '').replace(')', '').split(',')
    return tuple(map(int, parts))

def assemble_pptx(parent_metrics_dir, output_file):
    prs = Presentation()
    
    # 1. Set Slide Size (16:9 HD)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Clean Blank Slide Layout (Index 6)
    blank_slide_layout = prs.slide_layouts[6]

    # Find all slide subdirectories
    # Sort them to ensure order (slide_1, slide_2...)
    try:
        subdirs = sorted([
            d for d in os.listdir(parent_metrics_dir) 
            if os.path.isdir(os.path.join(parent_metrics_dir, d))
        ])
    except FileNotFoundError:
        print(f"Error: Directory {parent_metrics_dir} not found.")
        return

    print(f"Found {len(subdirs)} slides to assemble.")

    for slide_dir_name in subdirs:
        slide_path = os.path.join(parent_metrics_dir, slide_dir_name)
        layout_metrics_path = os.path.join(slide_path, 'layout_metrics.json')
        bg_image_path = os.path.join(slide_path, 'background.png')

        if not os.path.exists(layout_metrics_path):
            print(f"Skipping {slide_dir_name}: No metrics found.")
            continue

        print(f"Assembling slide: {slide_dir_name}")
        
        with open(layout_metrics_path, 'r') as f:
            metrics = json.load(f)

        # Add Slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # 2. Insert Background
        if os.path.exists(bg_image_path):
            slide.shapes.add_picture(bg_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Scale Factor
        SCALE_X = prs.slide_width / metrics['width']
        SCALE_Y = prs.slide_height / metrics['height']

        # 3. Create Placeholders
        for ph in metrics['placeholders']:
            rect = ph['rect']
            x = int(rect['x'] * SCALE_X)
            y = int(rect['y'] * SCALE_Y)
            w = int(rect['width'] * SCALE_X)
            h = int(rect['height'] * SCALE_Y)
            
            # Create Text Box
            textbox = slide.shapes.add_textbox(x, y, w, h)
            tf = textbox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = ph.get('text', '')
            
            # Apply Styles
            style = ph['style']
            
            # Font Size
            if style.get('fontSize'):
                px_size = float(style['fontSize'].replace('px', ''))
                p.font.size = Pt(px_size * 0.75)

            # Font Color
            if style.get('color'):
                try:
                    rgb = css_rgb_to_tuple(style['color'])
                    p.font.color.rgb = RGBColor(*rgb)
                except:
                    pass
            
            # Font Family
            if style.get('fontFamily'):
                p.font.name = style['fontFamily'].split(',')[0].strip().replace("'", "")

            
    
    prs.save(output_file)
    print(f"PPTX saved to: {output_file}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python assembler.py <metrics_dir_containing_json_and_bg> <output_pptx>")
        sys.exit(1)
        
    assemble_pptx(sys.argv[1], sys.argv[2])
